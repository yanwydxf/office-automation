import pptx
from pptx.util import Inches,Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import  RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
#步骤一：得到演示文稿的对象
prs=pptx.Presentation('test.pptx')
#步骤二：写入操作
slide=prs.slides.add_slide(prs.slide_layouts[0])
# prs.slides.add_slide(prs.slide_layouts[1])
# prs.slides.add_slide(prs.slide_layouts[2])
# #删除幻灯片
# print(len(prs.slides))
# del prs.slides._sldIdLst[1]
# print(len(prs.slides))
text1=slide.shapes.add_textbox(Inches(5),Inches(5),Inches(5),Inches(5))
text1.text="我是文本框"
p1=text1.text_frame.add_paragraph()
p1.text='我是段落1'
p1.add_run().text='end'
title_shape=slide.shapes.title
title_shape.text='标题1'
slide.shapes.placeholders[1].text='标题2'
#添加自选图形
shape=slide.shapes.add_shape(MSO_SHAPE.HEXAGON,Inches(2),Inches(2),Inches(5),Inches(3))
#填充、边框
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(255,0,0)
line=shape.line
line.color.rgb=RGBColor(55,3,5)
line.width=Pt(2)
#添加表格
table=slide.shapes.add_table(3,3,Inches(2),Inches(2),Inches(4),Inches(2)).table
#填充内容
table.cell(1,0).text='name'
table.cell(1,1).text='age'
table.cell(1,2).text='class'
table.cell(2,0).text='张三'
table.cell(2,1).text='19'
table.cell(2,2).text='一班'
#合并单元格
cell=table.cell(0,0)
cell1=table.cell(0,2)
cell.merge(cell1)
table.cell(0,0).text='班级学生信息' #第一行
print(cell.is_merge_origin)#单元格是否合并
cell.split()#取消合并
print(cell.is_merge_origin)
#写入图表
chart_data=CategoryChartData()
chart_data.categories=['一月份','二月份','三月份']#X轴
chart_data.add_series('Y2020',(300,400,500))
chart_data.add_series('Y2019',(400,500,900))
chart=slide.shapes.add_chart(XL_CHART_TYPE.LINE,Inches(2),Inches(2),Inches(6),Inches(4),chart_data).chart
chart.has_title=True
chart.chart_title.text_frame.text='第一季度销售额' #标题
chart.has_legend=True
chart.legend.position=XL_LEGEND_POSITION.RIGHT
#步骤三：保存PPT文件
prs.save('test.pptx')